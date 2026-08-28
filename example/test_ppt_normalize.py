# example/test_ppt_normalize.py
"""ppt_writer 共享契约层测试：spec 归一化规则 + manifest + 硬合规校验。

归一化覆盖：完整页按 page 值、缺省页按 section/模板、page 与 section 冲突
page 优先、缺 title 缺 section 兜底、非法输入含字段路径的错误信息。
硬合规覆盖：内置默认模板可验、合规/缺占位符/类型错误/不支持对象的逐项清单。
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from example.ppt_writer import normalize
from example.ppt_writer.normalize import (
    BUILTIN_LAYOUTS,
    KIND_REQUIREMENTS,
    normalize as norm,
    register_manifest,
    resolve_layout,
    verify_draft,
    verify_template,
    validate_manifest,
)


# ── 归一化：覆盖规则 ────────────────────────────────────────────────

def _spec(**over):
    spec = {
        "sections": [
            {"id": "intro", "title": "引言", "defaults": {"layout": "content"}},
            {"id": "end", "title": "致谢"},
        ],
        "pages": [
            {"index": 1, "section": "intro", "title": "完整页标题",
             "content": {"points": ["a", "b"]}, "layout": "content"},
            {"index": 2, "section": "end"},  # 缺省页：section 标题 + 模板兜底
            {"index": 3},                    # 无 section 无 title：全兜底
        ],
    }
    spec.update(over)
    return spec


def test_full_page_uses_page_values():
    pages = norm(_spec())
    p1 = pages[0]
    assert p1["title"] == "完整页标题"
    assert p1["layout"] == "content"
    assert p1["content"]["points"] == ["a", "b"]


def test_default_page_uses_section_title_and_template_layout():
    pages = norm(_spec())
    p2 = pages[1]
    assert p2["title"] == "致谢"          # section 标题
    assert p2["layout"] == "content"      # 模板兜底（section 无 defaults）
    assert p2["content"]["points"] == []  # 内容兜底为空 → 版式默认文本呈现
    assert p2["content"]["image"] is None


def test_page_wins_over_section_conflict():
    spec = _spec()
    spec["sections"][0]["defaults"]["layout"] = "thanks"
    spec["pages"][0]["title"] = "与 section 不一致的标题"
    p1 = norm(spec)[0]
    assert p1["title"] == "与 section 不一致的标题"  # page 优先，不报错
    assert p1["layout"] == "content"                  # page.layout 优先于 section 默认


def test_section_content_defaults_fill_page():
    spec = _spec()
    spec["sections"][0]["defaults"]["content"] = {
        "points": ["section 默认要点"], "caption": "section 默认注释",
    }
    p = norm(spec)[0]
    assert p["content"]["points"] == ["a", "b"]  # page 显式覆盖
    spec["sections"][1].setdefault("defaults", {})["content"] = {
        "points": ["section 默认要点"], "caption": "section 默认注释",
    }
    p2 = norm(spec)[1]  # 属于 section "end"
    assert p2["content"]["points"] == ["section 默认要点"]
    assert p2["content"]["caption"] == "section 默认注释"


def test_no_title_no_section_fallback():
    pages = norm(_spec())
    p3 = pages[2]
    assert p3["title"] == ""          # 无 page/section 标题 → 空 → 不填充
    assert p3["layout"] == "content"  # 模板兜底
    assert p3["section"] is None


def test_pages_sorted_by_index_and_levels_kept():
    spec = _spec()
    spec["pages"] = [
        {"index": 5, "title": "五"},
        {"index": 1, "title": "一"},
    ]
    pages = norm(spec)
    assert [p["index"] for p in pages] == [1, 5]
    spec["pages"][0]["content"] = {"points": [{"text": "子点", "level": 1}, "平点"]}
    pts = norm(spec)[1]["content"]["points"]  # 升序后 index 5 在位置 1
    assert pts == [{"text": "子点", "level": 1}, "平点"]


# ── 归一化：非法输入（含字段路径的错误信息）──────────────────────────

def test_invalid_spec_errors_carry_field_paths():
    cases = [
        ("[]", "顶层"),
        ({"pages": []}, "pages 必须是非空数组"),
        ({"pages": [{"title": "x"}]}, r"pages\[0\]\.index"),
        ({"pages": [{"index": "1"}]}, r"pages\[0\]\.index"),
        ({"pages": [{"index": 1, "section": "ghost"}]}, "section 'ghost'"),
        ({"sections": [{"id": "a"}], "pages": [{"index": 1}, {"index": 1}]},
         "index 重复"),
        ({"sections": [{"id": 1}], "pages": [{"index": 1}]}, r"sections\[0\]"),
        ({"pages": [{"index": 1, "content": {"points": [1]}}]}, r"points\[0\]"),
        ({"pages": [{"index": 1, "content": {"image": 5}}]}, "content.image"),
        ({"pages": [{"index": 1, "layout": "fancy"}]}, "layout"),
        ({"pages": [{"index": 1, "layout": "template:"}]}, "template:"),
    ]
    for spec, needle in cases:
        with pytest.raises(ValueError, match=needle):
            norm(json.loads(spec) if isinstance(spec, str) else spec)


# ── manifest：读写 / 解析 / 注册 ────────────────────────────────────

def test_manifest_missing_returns_empty(tmp_path: Path):
    assert normalize.load_manifest(tmp_path) == {}


def test_register_and_resolve(tmp_path: Path):
    register_manifest("content", {"file": "base.pptx", "layout": 3, "kind": "content"},
                      reference_dir=tmp_path)
    manifest = normalize.load_manifest(tmp_path)
    assert manifest["content"] == {"file": "base.pptx", "layout": 3, "kind": "content"}
    # 注册后覆盖
    register_manifest("content", {"file": "base.pptx", "layout": 4}, reference_dir=tmp_path)
    assert normalize.load_manifest(tmp_path)["content"]["layout"] == 4


def test_resolve_layout_builtin_and_manifest(tmp_path: Path):
    empty = normalize.load_manifest(tmp_path)  # 隔离：不读仓库真实 manifest
    assert resolve_layout("content", empty) == (None, 0, BUILTIN_LAYOUTS["content"], "content")
    assert resolve_layout("thanks", empty) == (None, 0, 5, "thanks")
    register_manifest("content", {"file": "mine.pptx", "layout": 2, "kind": "content"},
                      reference_dir=tmp_path)
    file, master, idx, kind = resolve_layout("content", normalize.load_manifest(tmp_path))
    assert (file, master, idx, kind) == ("mine.pptx", 0, 2, "content")
    register_manifest("alt", {"file": "mine.pptx", "master": 1, "layout": 2, "kind": "section"},
                      reference_dir=tmp_path)
    file, master, idx, kind = resolve_layout("alt", normalize.load_manifest(tmp_path))
    assert (file, master, idx, kind) == ("mine.pptx", 1, 2, "section")
    with pytest.raises(ValueError, match="未知逻辑模板 'nope'"):
        resolve_layout("nope", normalize.load_manifest(tmp_path))


def test_validate_manifest_catches_bad_entries():
    issues = validate_manifest({
        "a": "not-a-dict",
        "b": {"layout": "x"},
        "c": {"layout": 1, "kind": "nope"},
        "d": {"file": "f.pptx", "layout": 1, "master": -1},
    })
    text = "\n".join(issues)
    assert "manifest[a]" in text and "对象" in text
    assert "manifest[b]" in text and "layout" in text
    assert "manifest[c].kind" in text
    assert "manifest[d].master" in text


def test_verify_draft_master_index_bounds(tmp_path: Path):
    """单母版文件按 master=1 寻址 → 明确越界错误（多母版模板的寻址入口）。"""
    draft = _draft(tmp_path, 1)
    issues = verify_draft(draft, 1, "content", master_index=1)
    assert any("母版索引 1 越界" in i and "含 1 个母版" in i for i in issues)


def test_verify_draft_master_index_targets_second_master_layout(tmp_path: Path):
    """master=0 时仍按第一个母版寻址（缺省路径回归）。"""
    draft = _draft(tmp_path, 5)
    issues = verify_draft(draft, 5, "content", master_index=0)
    assert any("缺少文本占位符" in i for i in issues)


# ── 硬合规：内置默认模板全部可验 ────────────────────────────────────

def test_verify_template_builtin_defaults_pass():
    assert verify_template() == []  # 空 manifest → 内置 5 种 kind 全过


def test_verify_template_missing_file_reported(tmp_path: Path):
    register_manifest("content", {"file": "ghost.pptx", "layout": 0}, reference_dir=tmp_path)
    issues = verify_template(tmp_path)
    assert any("ghost.pptx" in i and "缺失" in i for i in issues)


def test_verify_template_bad_layout_index_reported(tmp_path: Path):
    from pptx import Presentation
    draft = tmp_path / "bad.pptx"
    Presentation().save(str(draft))
    register_manifest("content", {"file": "bad.pptx", "layout": 99}, reference_dir=tmp_path)
    issues = verify_template(tmp_path)
    assert any("越界" in i and "99" in i for i in issues)


# ── 硬合规：草稿（制作工作流判据）───────────────────────────────────

def _draft(tmp_path: Path, layout_index: int, renames: dict[str, str] | None = None) -> Path:
    from pptx import Presentation
    pres = Presentation()
    for ph in pres.slide_layouts[layout_index].placeholders:
        if ph.name in (renames or {}):
            ph.name = renames[ph.name]
    path = tmp_path / f"draft_{layout_index}.pptx"
    pres.save(str(path))
    return path


def test_verify_draft_compliant_content_and_picture(tmp_path: Path):
    assert verify_draft(_draft(tmp_path, 1), 1, "content") == []   # Title and Content
    assert verify_draft(_draft(tmp_path, 8), 8, "picture") == []   # Picture with Caption


def test_verify_draft_missing_placeholders_reported(tmp_path: Path):
    # Title Only（仅标题，无正文/图片占位符）按 content / picture 校验 → 逐项清单
    draft = _draft(tmp_path, 5)
    issues = verify_draft(draft, 5, "content")
    assert any("缺少文本占位符" in i for i in issues)
    issues = verify_draft(draft, 5, "picture")
    assert any("缺少文本占位符" in i for i in issues)
    assert any("缺少图片占位符" in i for i in issues)


def test_verify_draft_named_placeholders_pass_and_type_mismatch(tmp_path: Path):
    # 按名匹配：title/points 重命名后 content 合规
    draft = _draft(tmp_path, 1, {"Title 1": "title", "Content Placeholder 2": "points"})
    assert verify_draft(draft, 1, "content") == []
    # 名称匹配但类型不符：正文占位符命名为 image → 类型错误
    draft = _draft(tmp_path, 1, {"Content Placeholder 2": "image"})
    issues = verify_draft(draft, 1, "content")
    assert any("'image' 类型错误" in i for i in issues)


def test_verify_draft_missing_file():
    assert verify_draft("nope.pptx", 0, "content") == ["模板文件不存在: nope.pptx"]


# ── 模板形态诊断（页面层手工设计 / 未注册页面清单 / 修正指引）────────

def _hand_made_template(tmp_path: Path) -> Path:
    """自带 2 页：一页纯手工设计（空白版式 + 形状），一页占位符页（标题版式）。"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    pres = Presentation()
    s1 = pres.slides.add_slide(pres.slide_layouts[6])  # 空白版式
    s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 1000000, 1000000)
    s1.shapes.add_textbox(0, 0, 1000000, 1000000).text = "手工标题"
    pres.slides.add_slide(pres.slide_layouts[0])       # 标题幻灯片（占位符页）
    path = tmp_path / "hand_made.pptx"
    pres.save(str(path))
    return path


def test_diagnose_template_flags_hand_made_pages(tmp_path: Path):
    path = _hand_made_template(tmp_path)
    warnings = normalize.diagnose_template(path)
    text = "\n".join(warnings)
    assert "自带页 1 为手工设计" in text
    assert "0 占位符" in text
    assert "不会在渲染输出中复现" in text
    # 未注册页面清单（占位符页也列出）
    assert "自带页 2" in text and "占位符" in text
    # 修正指引
    assert "幻灯片母版" in text and "修正" in text


def test_diagnose_template_empty_and_missing(tmp_path: Path):
    from pptx import Presentation
    empty = tmp_path / "empty.pptx"
    Presentation().save(str(empty))
    assert normalize.diagnose_template(empty) == []  # 无自带页 → 无提醒
    assert any("不存在" in w for w in normalize.diagnose_template(tmp_path / "nope.pptx"))


def test_unsupported_object_detection(tmp_path: Path):
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("s", (1, 2))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, 0, 0, 3000000, 2000000, chart_data)
    assert "chart" in normalize._unsupported_objects(slide.shapes)


def test_kind_requirements_cover_all_builtin_kinds():
    for kind in BUILTIN_LAYOUTS:
        assert kind in KIND_REQUIREMENTS
