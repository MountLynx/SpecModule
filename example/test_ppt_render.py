# example/test_ppt_render.py
"""渲染读回验收测试（spec → .pptx → python-pptx 重开逐项断言）。

覆盖：3 页 spec（含图片页）渲染读回、页数/标题/要点/图片、中文文本、
未填充页显示版式默认文本（内置空默认 + 文件模板非空默认）、混合 deck
fail-fast、缺失图片报错。文件模板路径经 monkeypatch normalize.REFERENCE_DIR
指向临时目录（command 子进程路径由 run_render 冒烟测试覆盖）。
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from example.ppt_writer import normalize
from example.ppt_writer.module import (
    DEFAULT_SPEC,
    run_render,
    _build_registry,
)
from example.ppt_writer.render_deck import render as render_deck
from llm.mock import MockLLMClient

import base64

PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)  # 合法 1x1 PNG


def _spec_with_picture(tmp_path: Path, output: str) -> dict:
    img = tmp_path / "fig.png"
    img.write_bytes(PNG_BYTES)
    return {
        "title": "验收演示",
        "output": output,
        "theme": {"font": "微软雅黑"},
        "sections": [
            {"id": "intro", "title": "引言", "defaults": {"layout": "content"}},
            {"id": "end", "title": "致谢"},
        ],
        "pages": [
            {"index": 1, "section": "intro", "title": "验收标题", "layout": "title"},
            {
                "index": 2,
                "section": "intro",
                "title": "核心要点",
                "content": {"points": ["要点一：中文内容", "要点二", {"text": "子要点", "level": 1}]},
            },
            {
                "index": 3,
                "section": "intro",
                "title": "架构图",
                "layout": "picture",
                "content": {"image": str(img), "caption": "模块架构示意"},
            },
            {"index": 4, "section": "end"},  # 缺省页：标题取 section，内容走模板兜底
        ],
    }


def test_render_readback_full_flow(tmp_path: Path):
    """模板通道全链路：翻译器 → command 子进程渲染 → Report → 读回断言。"""
    import asyncio
    out = tmp_path / "deck.pptx"
    firings = asyncio.run(run_render(
        _spec_with_picture(tmp_path, str(out)), llm_client=MockLLMClient(),
        persist=False,
    ))
    report = firings[-1].output
    assert report["status"] == "ok", report
    assert report["pages"] == 4

    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    pres = Presentation(str(out))
    slides = list(pres.slides)
    assert len(slides) == 4

    # 每页标题与归一化后的 spec 一致（含中文）
    titles = []
    for s in slides:
        for ph in s.placeholders:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                titles.append(ph.text_frame.text)
    assert titles[0] == "验收标题"
    assert titles[1] == "核心要点"
    assert titles[2] == "架构图"
    assert titles[3] == "致谢"  # 缺省页标题取 section
    # 断言文本含中文（CJK 字符）
    assert any("\u4e00" <= c <= "\u9fff" for t in titles for c in t)

    # 要点逐项读回（Title and Content 的正文占位符类型为 OBJECT/BODY）
    body_texts = []
    for s in slides:
        for ph in s.placeholders:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                body_texts.append(ph.text_frame.text)
    assert "要点一：中文内容" in body_texts[0]
    assert "要点二" in body_texts[0]
    assert "子要点" in body_texts[0]

    # 图片页：图片对象存在（blip 嵌入占位符）+ caption 文本
    from pptx.oxml.ns import qn
    pic_slide = slides[2]
    assert sum(
        1 for s in pic_slide.shapes if len(list(s._element.iter(qn("a:blip")))) > 0
    ) == 1
    caption_texts = [
        ph.text_frame.text for ph in pic_slide.placeholders
        if ph.placeholder_format.type == PP_PLACEHOLDER.BODY  # Picture with Caption 的正文即 caption
    ]
    assert any("模块架构示意" in t for t in caption_texts)


def test_unfilled_page_shows_layout_default_text_in_file_template(tmp_path: Path, monkeypatch):
    """文件模板：未填充占位符保留版式默认文本（模板兜底规则在渲染期的实现）。"""
    # 制作带非空默认文本的模板：Title Only 版式标题占位符写入默认文案
    from pptx import Presentation
    ref_dir = tmp_path / "reference"
    ref_dir.mkdir()
    draft = ref_dir / "base.pptx"
    pres = Presentation()
    for ph in pres.slide_layouts[5].placeholders:  # Title Only
        ph.text_frame.text = "版式默认标题"
    pres.save(str(draft))
    normalize.register_manifest("thanks", {"file": "base.pptx", "layout": 5},
                                reference_dir=ref_dir)

    monkeypatch.setattr(normalize, "REFERENCE_DIR", ref_dir)
    out = tmp_path / "defaults.pptx"
    pages = normalize.normalize({
        "pages": [{"index": 1, "layout": "thanks"}],  # 无 title → 不填充，走文件模板
    })
    result = render_deck({
        "pages": pages, "output": str(out), "font": "微软雅黑",
    })
    assert result["pages"] == 1

    pres2 = Presentation(str(out))
    slide = list(pres2.slides)[0]
    # 未填充：幻灯片占位符文本为空（渲染器未触碰，布局默认文本保留）
    from pptx.enum.shapes import PP_PLACEHOLDER
    slide_titles = [
        ph.text_frame.text for ph in slide.placeholders
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    ]
    assert slide_titles == [""]
    # 版式默认文本留在模板（PowerPoint 渲染时呈现）——从版式读回验证
    layout_defaults = [
        ph.text_frame.text for ph in slide.slide_layout.placeholders
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    ]
    assert layout_defaults == ["版式默认标题"]


def test_mixed_file_and_builtin_deck_rejected(tmp_path: Path, monkeypatch):
    """文件 deck 中未注册的内置 kind → fail-fast 报缺注册（不清不楚的混合禁止）。"""
    from pptx import Presentation
    ref_dir = tmp_path / "reference"
    ref_dir.mkdir()
    draft = ref_dir / "base.pptx"
    Presentation().save(str(draft))
    normalize.register_manifest("content", {"file": "base.pptx", "layout": 1},
                                reference_dir=ref_dir)
    monkeypatch.setattr(normalize, "REFERENCE_DIR", ref_dir)
    pages = normalize.normalize({
        "pages": [
            {"index": 1, "title": "a", "layout": "template:content"},
            {"index": 2, "title": "b", "layout": "title"},  # 内置 kind 未注册
        ],
    })
    with pytest.raises(ValueError, match="'title' 未在 manifest 注册"):
        render_deck({"pages": pages, "output": str(tmp_path / "x.pptx"), "font": "微软雅黑"})


def test_file_template_sample_slides_cleared(tmp_path: Path, monkeypatch):
    """模板文件自带示例页时，输出只含 spec 页面（版式/主题复用，示例内容丢弃）。"""
    from pptx import Presentation
    ref_dir = tmp_path / "reference"
    ref_dir.mkdir()
    draft = ref_dir / "base.pptx"
    pres = Presentation()
    for _ in range(2):  # 模板自带 2 页示例
        pres.slides.add_slide(pres.slide_layouts[0])
    pres.save(str(draft))
    normalize.register_manifest("content", {"file": "base.pptx", "layout": 1},
                                reference_dir=ref_dir)
    normalize.register_manifest("title", {"file": "base.pptx", "layout": 0},
                                reference_dir=ref_dir)
    monkeypatch.setattr(normalize, "REFERENCE_DIR", ref_dir)

    out = tmp_path / "clean.pptx"
    pages = normalize.normalize({
        "pages": [
            {"index": 1, "title": "一", "layout": "title"},
            {"index": 2, "title": "二", "layout": "content", "content": {"points": ["p"]}},
        ],
    })
    result = render_deck({"pages": pages, "output": str(out), "font": "微软雅黑"})
    assert result["pages"] == 2
    assert len(list(Presentation(str(out)).slides)) == 2
    # 文件 deck 的形态诊断透出（未注册页面清单 + 修正指引分别成条）
    assert any("自带页" in w for w in result["warnings"])
    assert any("修正" in w and "幻灯片母版" in w for w in result["warnings"])


def test_missing_image_file_reported(tmp_path: Path):
    pages = normalize.normalize({
        "pages": [{"index": 1, "title": "a", "layout": "picture",
                   "content": {"image": str(tmp_path / "ghost.png")}}],
    })
    with pytest.raises(FileNotFoundError, match="ghost.png"):
        render_deck({"pages": pages, "output": str(tmp_path / "x.pptx"), "font": "微软雅黑"})


def test_default_spec_renders_three_pages(tmp_path: Path):
    """内置 3 页示例 spec 零配置可渲染。"""
    import asyncio
    out = tmp_path / "demo.pptx"
    spec = json.loads(json.dumps(DEFAULT_SPEC))
    spec["output"] = str(out)
    firings = asyncio.run(run_render(spec, llm_client=MockLLMClient(), persist=False))
    assert firings[-1].output["status"] == "ok"
    assert firings[-1].output["pages"] == 3


def test_registry_build_without_template_name():
    """build_registry 不依赖 template_name（双模板共享组件）。"""
    reg = _build_registry(MockLLMClient())
    assert reg.is_command("render_deck")
    assert reg.is_harness("template_opinion")
