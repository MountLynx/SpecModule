# example/ppt_writer/dump_template.py
"""模板 dump CLI：.pptx → 结构化描述 JSON（LLM 意见节点的输入）。

经制作工作流（template_review）的 dump script 调用：``python dump_template.py
--draft <path>``——每个版式输出占位符 index/name/类型/位置/字号/默认文本、
图片占位符有无、SmartArt/图表检测。意见 harness 只消费本 dump 的结构化描述
（pptx 是二进制，LLM 无法直接读）。

python-pptx 惰性导入；缺失时报含安装指引的清晰错误（exit 2）。
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

try:
    from . import normalize  # 包形态（模块引用/测试）
except ImportError:
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    import normalize as normalize  # 脚本直跑（command 子进程）

_EMU_PER_CM = 360000


def _ph_info(ph: Any) -> dict[str, Any]:
    """单个占位符的结构化描述。"""
    info: dict[str, Any] = {
        "index": ph.placeholder_format.idx,
        "name": ph.name,
        "type": normalize._placeholder_type_name(ph) or None,
        "x_cm": round((ph.left or 0) / _EMU_PER_CM, 2),
        "y_cm": round((ph.top or 0) / _EMU_PER_CM, 2),
        "w_cm": round((ph.width or 0) / _EMU_PER_CM, 2),
        "h_cm": round((ph.height or 0) / _EMU_PER_CM, 2),
        "default_text": ph.text_frame.text if ph.has_text_frame else "",
    }
    if ph.has_text_frame:
        runs = [r for p in ph.text_frame.paragraphs for r in p.runs]
        size = runs[0].font.size if runs else None
        info["font_size"] = round(size.pt, 1) if size is not None else None
    return info


def dump_pptx(draft_path: Path | str) -> dict[str, Any]:
    """.pptx 草稿 → 结构化描述 JSON（每母版每版式占位符/图片/不支持对象）。

    覆盖**全部母版**（多母版模板的变体版式按母版索引区分——渲染/manifest
    寻址与 dump 一致）。
    """
    normalize._ensure_pptx()
    from pptx import Presentation  # noqa: PLC0415
    draft_path = Path(draft_path)
    if not draft_path.exists():
        raise FileNotFoundError(f"草稿文件不存在: {draft_path}")
    pres = Presentation(str(draft_path))
    masters: list[dict[str, Any]] = []
    for mi, master in enumerate(pres.slide_masters):
        layouts: list[dict[str, Any]] = []
        for i, layout in enumerate(master.slide_layouts):
            placeholders = [
                _ph_info(ph) for ph in layout.placeholders
            ]
            unsupported = normalize._unsupported_objects(layout.shapes)
            layouts.append({
                "index": i,
                "name": layout.name,
                "placeholders": placeholders,
                "has_picture_placeholder": any(
                    p["type"] == "PICTURE" for p in placeholders
                ),
                "unsupported": unsupported,
            })
        masters.append({"index": mi, "layouts": layouts})
    return {"file": str(draft_path), "masters": masters}


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="ppt_writer 模板 dump（LLM 意见输入）")
    parser.add_argument("--draft", required=True, help=".pptx 草稿路径")
    args = parser.parse_args(argv)
    try:
        normalize._ensure_pptx()
    except ImportError as e:
        print(f"错误: {e}", file=sys.stderr)
        return 2
    try:
        dump = dump_pptx(args.draft)
    except (ValueError, FileNotFoundError, RuntimeError) as e:
        print(f"dump 失败: {e}", file=sys.stderr)
        return 1
    print(json.dumps(dump, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
