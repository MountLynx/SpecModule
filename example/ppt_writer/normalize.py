# example/ppt_writer/normalize.py
"""共享契约模块：manifest 读写 + 逻辑模板解析 + 硬合规校验 + spec 归一化。

本文件是 example/ppt_writer 的**共享契约层**（AGENTS.md：确认第二消费者才抽——
渲染器 render_deck / 制作工作流 template_review / 测试三端消费）：

- ``manifest``：逻辑模板名 → {file, layout, kind}（file 省略 = 内置默认模板）
- ``resolve_layout``：逻辑名 → (文件 | None, 版式索引)
- 模板/参考目录一律按**本模块文件相对路径**定位（command 节点 cwd 不定，
  从哪启动不影响结果）
- ``verify_*``：硬合规校验纯函数（渲染前 fail-fast 与制作工作流共用）
- ``normalize``：spec → pages_final（page 显式 > section 默认 > 模板兜底，
  冲突消解不报错；非法 spec 抛出含字段路径的明确错误）

python-pptx 一律**惰性导入**（``_ensure_pptx``）——模块加载不依赖 pptx，
缺失时报含安装指引的清晰错误。
"""

from __future__ import annotations

import json
import os
import tempfile
from pathlib import Path
from typing import Any

# ── 模块文件相对路径定位（不依赖 cwd）───────────────────────────────
MODULE_DIR = Path(__file__).resolve().parent
REFERENCE_DIR = MODULE_DIR / "reference"
MANIFEST_PATH = REFERENCE_DIR / "manifest.json"

# 内置默认模板：python-pptx 内置 11 版式映射（无 manifest 条目时兜底）。
# Title Slide(0) / Title and Content(1) / Section Header(2) / Title Only(5)
# / Picture with Caption(8)
BUILTIN_LAYOUTS: dict[str, int] = {
    "title": 0,
    "section": 2,
    "content": 1,
    "picture": 8,
    "thanks": 5,
}
BUILTIN_KINDS = tuple(BUILTIN_LAYOUTS)

# 每种版式 kind 渲染所需占位符能力（硬合规判据）：
# - title/section/thanks：仅标题
# - content：标题 + 至少一个文本占位符（要点）
# - picture：标题 + 图片占位符 + 至少一个文本占位符（语义注释 caption）
KIND_REQUIREMENTS: dict[str, list[str]] = {
    "title": ["title"],
    "section": ["title"],
    "thanks": ["title"],
    "content": ["title", "text"],
    "picture": ["title", "image", "text"],
}

# 占位符命名约定（模板作者在 PowerPoint 里重命名）；按名匹配优先，
# 未命名时按占位符类型兜底（内置版式即走兜底路径）。
ROLE_NAMES = ("title", "points", "caption", "image")

DEFAULT_FONT = "微软雅黑"

_ENVELOPE_DIR = Path(tempfile.gettempdir())


# ── python-pptx 惰性导入 ────────────────────────────────────────────

def _ensure_pptx() -> Any:
    """惰性导入 python-pptx；缺失时抛含安装指引的清晰错误。"""
    try:
        import pptx  # noqa: PLC0415
        return pptx
    except ImportError as e:
        raise ImportError(
            "缺少 python-pptx——example/ppt_writer 的渲染/模板工作流需要它。"
            "安装：pip install python-pptx（或 pip install -r requirements.txt）"
        ) from e


def _pptx_enums(pptx: Any) -> tuple[Any, Any]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER  # noqa: PLC0415
    return MSO_SHAPE_TYPE, PP_PLACEHOLDER


# ── manifest：读写 / 校验 / 注册 ────────────────────────────────────

def manifest_path(reference_dir: Path | str | None = None) -> Path:
    """manifest 文件路径（缺省 = 模块 reference/）。"""
    return (Path(reference_dir) if reference_dir is not None else REFERENCE_DIR) / "manifest.json"


def load_manifest(reference_dir: Path | str | None = None) -> dict[str, Any]:
    """读取 manifest：{逻辑模板名: {file?, layout, kind?}}。缺失/损坏报明确错误。"""
    path = manifest_path(reference_dir)
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as e:
        raise ValueError(f"manifest 读取失败（{path}）: {e}") from e
    if not isinstance(data, dict):
        raise ValueError(f"manifest 顶层必须是对象（{path}）")
    return data


def validate_manifest(manifest: dict[str, Any]) -> list[str]:
    """manifest 结构校验：逐条目检查字段类型与取值合法性（不含文件存在性）。"""
    issues: list[str] = []
    for name, entry in manifest.items():
        if not isinstance(entry, dict):
            issues.append(f"manifest[{name}]: 条目必须是对象")
            continue
        if "layout" not in entry:
            issues.append(f"manifest[{name}]: 缺少 'layout'（版式索引）")
        elif not isinstance(entry["layout"], int) or isinstance(entry["layout"], bool):
            issues.append(f"manifest[{name}].layout: 必须是整数版式索引")
        if "master" in entry and (
            not isinstance(entry["master"], int) or isinstance(entry["master"], bool)
            or entry["master"] < 0
        ):
            issues.append(f"manifest[{name}].master: 必须是非负整数母版索引（缺省 0 = 第一个母版）")
        if "file" in entry and not isinstance(entry["file"], str):
            issues.append(f"manifest[{name}].file: 必须是字符串（相对 reference/ 的文件名）")
        if "kind" in entry:
            if entry["kind"] not in KIND_REQUIREMENTS:
                issues.append(
                    f"manifest[{name}].kind: 未知 kind {entry['kind']!r}"
                    f"（可选: {', '.join(KIND_REQUIREMENTS)}）"
                )
        elif name in BUILTIN_KINDS:
            # 内置名未声明 kind → 按名推断
            pass
    return issues


def register_manifest(
    name: str,
    entry: dict[str, Any],
    reference_dir: Path | str | None = None,
) -> dict[str, Any]:
    """注册/覆盖一条 manifest 条目（file 相对 reference/）。返回更新后的 manifest。"""
    path = manifest_path(reference_dir)
    manifest = load_manifest(reference_dir)
    manifest[name] = entry
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return manifest


def resolve_layout(
    name: str,
    manifest: dict[str, Any] | None = None,
) -> tuple[str | None, int, int, str]:
    """逻辑模板名 → (文件相对名 | None, 母版索引, 版式索引, kind)。

    manifest 有条目 → 按条目（file 省略 = 内置默认模板文件；master 缺省 0）；
    无条目 → 内置默认映射（仅限 title/section/content/picture/thanks，
    master 0）。
    未知逻辑名 → ValueError（含可用清单）。
    """
    manifest = load_manifest() if manifest is None else manifest
    entry = manifest.get(name)
    if entry is not None:
        if not isinstance(entry, dict) or "layout" not in entry:
            raise ValueError(f"manifest[{name}]: 条目无效（缺少 layout）")
        kind = entry.get("kind") or (name if name in BUILTIN_KINDS else "content")
        master = entry.get("master", 0)
        return entry.get("file"), master, entry["layout"], kind
    if name in BUILTIN_LAYOUTS:
        return None, 0, BUILTIN_LAYOUTS[name], name
    available = ", ".join(sorted(set(BUILTIN_KINDS) | set(manifest)))
    raise ValueError(
        f"未知逻辑模板 '{name}'——manifest 无此条目且非内置 kind。"
        f"可用: {available}；自定义模板请用 layout: 'template:{name}' 并在 manifest 注册"
    )


# ── 占位符角色解析（渲染器与校验共用）───────────────────────────────

def _placeholder_type_name(ph: Any) -> str:
    ptype = ph.placeholder_format.type
    return getattr(ptype, "name", str(ptype)) if ptype is not None else ""


def _type_role(ptype_name: str) -> str | None:
    """占位符类型 → 角色（未命名时的兜底映射）。"""
    if ptype_name in ("TITLE", "CENTER_TITLE", "SUBTITLE"):
        return "title" if ptype_name != "SUBTITLE" else "caption"
    if ptype_name == "PICTURE":
        return "image"
    if ptype_name in ("BODY", "OBJECT", "TEXT"):
        return "points" if ptype_name in ("BODY", "OBJECT") else "caption"
    return None  # DATE / FOOTER / SLIDE_NUMBER / TABLE / CHART ... 忽略


_TYPE_EXPECT: dict[str, tuple[str, ...]] = {
    "title": ("TITLE", "CENTER_TITLE"),
    "points": ("BODY", "OBJECT", "TEXT"),
    "caption": ("BODY", "OBJECT", "TEXT"),
    "image": ("PICTURE",),
}


def _resolve_roles(placeholders: list[Any]) -> tuple[dict[str, Any], list[str]]:
    """占位符 → 角色映射（渲染填充与硬合规共用同一套解析）。

    规则：按名（title/points/caption/image）匹配优先，且类型须与角色相符；
    名称匹配但类型不符 → 记问题且不采用（回落类型兜底）；未命名占位符按
    类型兜底（TITLE→title、PICTURE→image、BODY/OBJECT→points、
    TEXT/SUBTITLE→caption、剩余文本占位符亦可作 caption）。同一占位符
    不重复分配角色。返回 (roles, 逐项类型问题清单)。
    """
    roles: dict[str, Any] = {}
    issues: list[str] = []
    assigned: set[int] = set()

    def try_assign(role: str, ph: Any) -> None:
        if id(ph) not in assigned and role not in roles:
            roles[role] = ph
            assigned.add(id(ph))

    # pass 1：按名
    for ph in placeholders:
        name = (ph.name or "").strip().lower()
        if name not in ROLE_NAMES:
            continue
        ptype = _placeholder_type_name(ph)
        if ptype in _TYPE_EXPECT[name]:
            try_assign(name, ph)
        else:
            issues.append(
                f"占位符 '{ph.name}' 类型错误：应为 {_TYPE_EXPECT[name][0]} 类型，"
                f"实际为 {ptype or '无类型'}"
            )
    # pass 2：按类型兜底
    for ph in placeholders:
        if id(ph) in assigned:
            continue
        role = _type_role(_placeholder_type_name(ph))
        if role is not None:
            try_assign(role, ph)
    # caption 兜底：未被 title/points/image 占用的文本占位符
    text_phs = [
        ph for ph in placeholders
        if id(ph) not in assigned and _type_role(_placeholder_type_name(ph)) in (
            "points", "caption"
        )
    ]
    if "caption" not in roles and text_phs:
        roles["caption"] = text_phs[0]
    return roles, issues


def _text_capable(ph: Any) -> bool:
    return _type_role(_placeholder_type_name(ph)) in ("title", "points", "caption")


# ── 硬合规校验（纯函数，渲染器与制作工作流共用）──────────────────────

def _unsupported_objects(shapes: Any) -> list[str]:
    """检测版式中的 v1 不支持对象（SmartArt / 图表 / 表格）。"""
    found: list[str] = []
    from pptx.oxml.ns import qn  # noqa: PLC0415
    for shape in shapes:
        if shape.shape_type == 3:  # MSO_SHAPE_TYPE.CHART
            found.append("chart")
            continue
        if shape.has_table:  # type: ignore[attr-defined]
            found.append("table")
            continue
        gd = shape._element.find(qn("a:graphicData"))
        if gd is None:
            continue
        uri = gd.get("uri", "")
        if "diagram" in uri:
            found.append("smartart")
        elif "chart" in uri:
            found.append("chart")
        elif "graphicFrame" in uri:
            found.append("table")
    return sorted(set(found))


def verify_layout(layout: Any, kind: str) -> list[str]:
    """校验单个版式是否满足 kind 的渲染需求。返回逐项问题清单（空 = 合规）。"""
    issues: list[str] = []
    placeholders = list(layout.placeholders)
    roles, name_issues = _resolve_roles(placeholders)
    issues.extend(name_issues)
    if "title" not in roles:
        issues.append("缺少标题占位符（命名为 title，或 TITLE 类型）")
    req = KIND_REQUIREMENTS.get(kind)
    if req is None:
        issues.append(f"未知 kind {kind!r}（可选: {', '.join(KIND_REQUIREMENTS)}）")
        req = []
    if "text" in req:
        texts = [ph for ph in placeholders if _text_capable(ph) and roles.get("title") is not ph]
        if not texts:
            issues.append("缺少文本占位符（points/caption——要点或图片语义注释）")
    if "image" in req and "image" not in roles:
        issues.append("缺少图片占位符（命名为 image，或 PICTURE 类型）")
    unsupported = _unsupported_objects(layout.shapes)
    if unsupported:
        issues.append(f"含 v1 不支持的填充对象: {', '.join(unsupported)}")
    return issues


def _open_layout(
    file_path: Path,
    layout_index: int,
    master_index: int = 0,
) -> tuple[Any, str | None]:
    """打开模板文件返回 (版式, 错误信息)；母版/版式索引越界报明确错误。"""
    pptx = _ensure_pptx()
    if not file_path.exists():
        return None, f"模板文件不存在: {file_path}"
    try:
        pres = pptx.Presentation(str(file_path))
    except Exception as e:
        return None, f"模板文件无法打开（{file_path}）: {e}"
    masters = pres.slide_masters
    if master_index < 0 or master_index >= len(masters):
        return None, (
            f"母版索引 {master_index} 越界——文件含 {len(masters)} 个母版（索引 0..{len(masters) - 1}）"
        )
    layouts = masters[master_index].slide_layouts
    if layout_index < 0 or layout_index >= len(layouts):
        return None, (
            f"版式索引 {layout_index} 越界——母版 {master_index} 含 {len(layouts)} 个版式"
            f"（索引 0..{len(layouts) - 1}）"
        )
    return layouts[layout_index], None


def verify_template(
    reference_dir: Path | str | None = None,
    manifest: dict[str, Any] | None = None,
) -> list[str]:
    """硬合规校验整个模板池：manifest 逻辑模板齐全 + 各版式占位符/类型 +
    无 v1 不支持对象。渲染器与制作工作流共用（fail-fast）。

    覆盖：
    - manifest 结构问题（validate_manifest）
    - 每条目引用的文件存在、版式索引有效、占位符满足 kind 需求、无 SmartArt/
      图表/表格
    - 内置默认模板（未在 manifest 覆盖的 kind）同样逐版式校验
    返回逐项问题清单（空 = 合规）。
    """
    manifest = load_manifest(reference_dir) if manifest is None else manifest
    issues: list[str] = validate_manifest(manifest)
    ref_dir = Path(reference_dir) if reference_dir is not None else REFERENCE_DIR

    # manifest 条目（文件模板）
    for name, entry in manifest.items():
        if not isinstance(entry, dict):
            continue
        prefix = f"manifest[{name}]"
        if "file" not in entry:
            continue  # 纯内置别名，无需文件校验
        file_path = ref_dir / entry["file"]
        if not file_path.exists():
            issues.append(f"{prefix}.file 缺失: {file_path}")
            continue
        layout, err = _open_layout(file_path, entry["layout"], entry.get("master", 0))
        if err is not None:
            issues.append(f"{prefix}: {err}")
            continue
        kind = entry.get("kind") or (name if name in BUILTIN_KINDS else "content")
        for issue in verify_layout(layout, kind):
            issues.append(f"{prefix}: {issue}")

    # 内置默认模板（manifest 未覆盖的 kind）
    for kind, idx in BUILTIN_LAYOUTS.items():
        if kind in manifest and isinstance(manifest[kind], dict) and "file" in manifest[kind]:
            continue
        pres = _ensure_pptx().Presentation()
        layout = pres.slide_layouts[idx]
        for issue in verify_layout(layout, kind):
            issues.append(f"内置默认模板 {kind}（版式 {idx}）: {issue}")
    return issues


def verify_draft(
    draft_path: Path | str,
    layout_index: int,
    kind: str,
    master_index: int = 0,
) -> list[str]:
    """制作工作流硬合规：校验草稿文件指定母版/版式（未入库前，渲染器同判据）。"""
    layout, err = _open_layout(Path(draft_path), layout_index, master_index)
    if err is not None:
        return [err]
    return [
        f"草稿母版 {master_index} 版式 {layout_index}: {issue}"
        for issue in verify_layout(layout, kind)
    ]


# ── 模板形态诊断（advisory，不阻断硬合规 gate）──────────────────────

def diagnose_template(draft_path: Path | str) -> list[str]:
    """诊断模板文件形态，输出逐项提醒（空 = 无提醒）。

    覆盖两类 v1 语义外的情况（审查时明确告知，渲染时随 warnings 透出）：
    1. **页面层手工设计**：自带 slide 无占位符、全为手工绘制的形状/图片/
       文本框——设计在页面上，版式层是默认样式；渲染只复用版式/主题，
       这些页面样式**不会复现**；
    2. **未注册页面清单**：自带页使用的版式/母版逐条列出——它们不在
       manifest 注册范围内，不会出现在渲染输出。
    附简略修正方法（把设计迁移到版式层后重新审查入库）。
    """
    if not Path(draft_path).exists():
        return [f"模板文件不存在: {draft_path}"]
    try:
        pres = _ensure_pptx().Presentation(str(draft_path))
    except Exception as e:
        return [f"模板文件无法打开（{draft_path}）: {e}"]
    masters = list(pres.slide_masters)
    slides = list(pres.slides)
    if not slides:
        return []
    warnings: list[str] = []
    hand_made = 0
    for si, slide in enumerate(slides, 1):
        phs = [s for s in slide.shapes if s.is_placeholder]
        extras = [s for s in slide.shapes if not s.is_placeholder]
        layout = slide.slide_layout
        mi = masters.index(layout.slide_master)
        kinds = [str(s.shape_type).split()[0] for s in extras]
        if not phs and extras:
            hand_made += 1
            warnings.append(
                f"自带页 {si} 为手工设计：0 占位符、{len(extras)} 个非占位符对象"
                f"（{', '.join(dict.fromkeys(kinds))}，版式「{layout.name}」/母版 {mi}）"
                f"——该页样式不会在渲染输出中复现"
            )
        else:
            warnings.append(
                f"自带页 {si}：{len(phs)} 个占位符 + {len(extras)} 个非占位符对象"
                f"（版式「{layout.name}」/母版 {mi}）——页面不在 manifest 注册范围，"
                f"不会出现在渲染输出"
            )
    if hand_made == len(slides):
        warnings.append(
            "模板的视觉设计全部位于页面层（版式层为 Office 默认样式）。"
            "渲染只复用版式/主题，输出将是默认外观。"
        )
    warnings.append(
        "简略修正方法：在 PowerPoint「视图 → 幻灯片母版」中把页面设计迁移到版式层"
        "（背景/配色/装饰对象 + 占位符排版），保存后重新运行本审查入库；"
        "此后渲染即可复现样式。"
    )
    return warnings


# ── 渲染前 fail-fast：按实际用到的版式解析 + 校验 ───────────────────

def resolve_deck_layouts(
    used_layouts: list[str],
    reference_dir: Path | str | None = None,
    manifest: dict[str, Any] | None = None,
) -> tuple[str | None, dict[str, tuple[int, int, str]], list[str]]:
    """把页面使用的逻辑模板名解析为渲染计划。

    返回 (deck_file | None, {逻辑名: (母版索引, 版式索引, kind)}, 问题清单)。

    - 全部内置 kind 且无文件条目 → deck_file=None（内置默认模板）
    - 任一逻辑名落到文件模板 → deck 文件 = 该文件；其余逻辑名必须指向
      **同一文件**（一个 Presentation 只能加载一套模板文件），否则报错
    - 文件 deck 中未在 manifest 覆盖的内置 kind → 报错（清单齐全性检查）
    - manifest 条目可含可选 ``master``（缺省 0）——多母版模板（如一套主题
      多个变体）按母版索引寻址版式
    """
    manifest = load_manifest(reference_dir) if manifest is None else manifest
    issues: list[str] = list(validate_manifest(manifest))
    plan: dict[str, tuple[int, int, str]] = {}
    deck_file: str | None = None

    for name in dict.fromkeys(used_layouts):  # 去重保序
        lookup = name[len("template:"):] if name.startswith("template:") else name
        entry = manifest.get(lookup)
        if entry is not None:
            if not isinstance(entry, dict) or "file" not in entry or "layout" not in entry:
                issues.append(
                    f"逻辑模板 '{lookup}' 的 manifest 条目无效"
                    f"（template: 引用必须有 file + layout）"
                )
                continue
            plan[name] = (
                entry.get("master", 0),
                entry["layout"],
                entry.get("kind") or (
                    lookup if lookup in BUILTIN_KINDS else "content"
                ),
            )
            if deck_file is None:
                deck_file = entry["file"]
            elif entry["file"] != deck_file:
                issues.append(
                    f"逻辑模板 '{lookup}' 与 '{deck_file}' 不在同一模板文件"
                    f"（一个 deck 只能使用一个模板文件）"
                )
            continue
        if name in BUILTIN_LAYOUTS:
            if deck_file is not None:
                issues.append(
                    f"逻辑模板 '{name}' 未在 manifest 注册（deck 为文件模板 "
                    f"'{deck_file}'，内置 kind 需逐一映射到该文件版式）"
                )
                continue
            plan[name] = (0, BUILTIN_LAYOUTS[name], name)
            continue
        issues.append(f"未知逻辑模板 '{name}'（未在 manifest 注册且非内置 kind）")

    # 文件模板校验（fail-fast 与 verify_template 同判据）
    ref_dir = Path(reference_dir) if reference_dir is not None else REFERENCE_DIR
    if deck_file is not None:
        file_path = ref_dir / deck_file
        if not file_path.exists():
            issues.append(f"模板文件缺失: {file_path}")
            return deck_file, plan, issues
        for name, (m_idx, idx, kind) in plan.items():
            layout, err = _open_layout(file_path, idx, m_idx)
            if err is not None:
                issues.append(f"manifest[{name}]: {err}")
                continue
            for issue in verify_layout(layout, kind):
                issues.append(f"manifest[{name}]: {issue}")
    else:
        pres = _ensure_pptx().Presentation()
        for name, (m_idx, idx, kind) in plan.items():
            for issue in verify_layout(pres.slide_layouts[idx], kind):
                issues.append(f"内置默认模板 {name}（版式 {idx}）: {issue}")
    return deck_file, plan, issues


# ── spec 归一化（纯函数）────────────────────────────────────────────

def validate_spec(spec: Any) -> None:
    """spec 结构校验；非法输入抛出含字段路径的明确错误。"""
    if not isinstance(spec, dict):
        raise ValueError(f"非法 spec: 顶层必须是 JSON 对象，实际 {type(spec).__name__}")
    sections = spec.get("sections", [])
    if not isinstance(sections, list):
        raise ValueError("非法 spec: sections 必须是数组")
    section_ids: set[str] = set()
    for i, sec in enumerate(sections):
        if not isinstance(sec, dict):
            raise ValueError(f"非法 spec: sections[{i}] 必须是对象")
        if "id" not in sec or not isinstance(sec["id"], str) or not sec["id"]:
            raise ValueError(f"非法 spec: sections[{i}] 缺少非空字符串 id")
        section_ids.add(sec["id"])
        defaults = sec.get("defaults")
        if defaults is not None and not isinstance(defaults, dict):
            raise ValueError(f"非法 spec: sections[{i}].defaults 必须是对象")
        if defaults is not None and "layout" in defaults and not isinstance(defaults["layout"], str):
            raise ValueError(f"非法 spec: sections[{i}].defaults.layout 必须是字符串")
        if defaults is not None and "content" in defaults:
            _validate_content(f"sections[{i}].defaults.content", defaults["content"])
    pages = spec.get("pages")
    if not isinstance(pages, list) or not pages:
        raise ValueError("非法 spec: pages 必须是非空数组")
    seen_idx: set[int] = set()
    for i, page in enumerate(pages):
        if not isinstance(page, dict):
            raise ValueError(f"非法 spec: pages[{i}] 必须是对象")
        if "index" not in page or not isinstance(page["index"], int) or isinstance(page["index"], bool):
            raise ValueError(f"非法 spec: pages[{i}].index 缺失或不是整数")
        if page["index"] in seen_idx:
            raise ValueError(f"非法 spec: pages[{i}].index 重复（{page['index']}）")
        seen_idx.add(page["index"])
        if "section" in page:
            if not isinstance(page["section"], str):
                raise ValueError(f"非法 spec: pages[{i}].section 必须是字符串")
            if page["section"] not in section_ids:
                raise ValueError(
                    f"非法 spec: pages[{i}].section '{page['section']}' 未在 sections 中定义"
                )
        if "title" in page and not isinstance(page["title"], str):
            raise ValueError(f"非法 spec: pages[{i}].title 必须是字符串")
        if "layout" in page:
            _validate_layout(f"pages[{i}].layout", page["layout"])
        if "content" in page:
            _validate_content(f"pages[{i}].content", page["content"])


def _validate_layout(field: str, layout: Any) -> None:
    if not isinstance(layout, str):
        raise ValueError(f"非法 spec: {field} 必须是字符串（可选: {', '.join(BUILTIN_KINDS)} 或 template:<名>）")
    if layout in BUILTIN_KINDS:
        return
    if layout.startswith("template:"):
        name = layout[len("template:"):]
        if not name:
            raise ValueError(f"非法 spec: {field} 'template:' 后必须跟逻辑模板名")
        return
    raise ValueError(
        f"非法 spec: {field} 未知版式 '{layout}'"
        f"（可选: {', '.join(BUILTIN_KINDS)} 或 template:<manifest 逻辑名>）"
    )


def _validate_content(field: str, content: Any) -> None:
    if not isinstance(content, dict):
        raise ValueError(f"非法 spec: {field} 必须是对象")
    if "points" in content:
        if not isinstance(content["points"], list):
            raise ValueError(f"非法 spec: {field}.points 必须是数组")
        for j, pt in enumerate(content["points"]):
            if isinstance(pt, str):
                continue
            if isinstance(pt, dict) and isinstance(pt.get("text"), str) and (
                "level" not in pt or (isinstance(pt["level"], int) and not isinstance(pt["level"], bool))
            ):
                continue
            raise ValueError(
                f"非法 spec: {field}.points[{j}] 必须是字符串或 {{text: str, level?: int}}"
            )
    for key in ("image", "caption"):
        if key in content and not isinstance(content[key], str):
            raise ValueError(f"非法 spec: {field}.{key} 必须是字符串")


def normalize(spec: dict[str, Any]) -> list[dict[str, Any]]:
    """spec → pages_final（每页定稿 title/content/layout）。

    覆盖规则（从最小一级向上，冲突消解不报错）：
    - title：page.title > section.title > ''（空 → 渲染时不填，版式默认文本兜底）
    - layout：page.layout > section.defaults.layout > 'content'
    - content：page.content 逐字段覆盖 section.defaults.content > 模板兜底
      （points=[] / image=None / caption=None）

    返回按 index 升序的页面列表：
    [{index, section?, title, content: {points, image, caption}, layout}]
    """
    validate_spec(spec)
    sections = {sec["id"]: sec for sec in spec.get("sections", [])}
    pages_final: list[dict[str, Any]] = []
    for page in spec["pages"]:
        section = sections.get(page["section"]) if page.get("section") else None
        defaults = (section or {}).get("defaults") or {}
        sec_content = defaults.get("content") or {}
        page_content = page.get("content") or {}
        title = page.get("title")
        if title is None:
            title = (section or {}).get("title", "")
        layout = page.get("layout")
        if layout is None:
            layout = defaults.get("layout", "content")
        points = page_content.get("points", sec_content.get("points", []))
        image = page_content.get("image", sec_content.get("image"))
        caption = page_content.get("caption", sec_content.get("caption"))
        pages_final.append({
            "index": page["index"],
            "section": page.get("section"),
            "title": title if isinstance(title, str) else "",
            "content": {
                "points": [dict(p) if isinstance(p, dict) else p for p in points],
                "image": image,
                "caption": caption,
            },
            "layout": layout,
        })
    pages_final.sort(key=lambda p: p["index"])
    return pages_final


# ── 运行信封（translator → command/script 的数据通道）────────────────
# command 节点的命令字符串是静态的（框架约束），spec 数据无法经节点输入
# 传给 command——翻译器（模板通道，能看到 spec）把归一化结果写入一个
# 进程级临时信封文件，command/script 节点按同路径读取（同一进程内注册
# 与执行，路径稳定）。

def _envelope_path(kind: str) -> Path:
    return _ENVELOPE_DIR / f"specmodule_ppt_writer_{kind}_{os.getpid()}.json"


def write_envelope(kind: str, data: dict[str, Any]) -> Path:
    """写入进程级信封（translator 调用）。返回路径。"""
    path = _envelope_path(kind)
    path.write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
    return path


def read_envelope(kind: str) -> dict[str, Any]:
    """读取进程级信封（command/script 节点调用）。缺失报明确错误。"""
    return read_envelope_file(_envelope_path(kind))


def read_envelope_file(path: Path | str) -> dict[str, Any]:
    """按路径读取信封 JSON（command 子进程用 --envelope 传入的路径）。"""
    path = Path(path)
    if not path.exists():
        raise RuntimeError(
            f"缺少运行信封 {path}——模板通道未先执行翻译器？"
            "（ppt_render/template_review 必须经模板通道运行，不能直接 --tasklist）"
        )
    return json.loads(path.read_text(encoding="utf-8"))
