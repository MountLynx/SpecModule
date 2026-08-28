# example/ppt_writer/render_deck.py
"""渲染器 CLI：spec 归一化结果（信封）→ .pptx。

经 command 节点以子进程调用（``python render_deck.py --envelope <path>``）：
读信封（翻译器写入的 pages_final + 输出路径 + 字体）→ 渲染前 fail-fast
硬合规（复用 normalize 的 verify）→ python-pptx 逐页渲染 → 保存。

零 LLM。python-pptx 惰性导入，缺失时报含安装指引的清晰错误（exit 2）。
成功：stdout 输出 JSON {pptx, pages, warnings}；失败：stderr 输出问题清单，
exit 1（合规问题）或 2（环境缺失）。
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

try:
    from . import normalize  # 包形态（模块引用/测试）
except ImportError:
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    import normalize as normalize  # 脚本直跑（command 子进程）


def _set_run_font(run: Any, font_name: str) -> None:
    """显式设置中文字体（latin + eastAsia 双 typeface，Windows 上避免回退替换）。"""
    run.font.name = font_name
    from pptx.oxml.ns import qn  # noqa: PLC0415
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font_name)


def _fill_text_paragraph(tf: Any, text: str, font_name: str, level: int = 0) -> None:
    tf.clear()  # 清掉版式默认文本，仅留一个空段落
    p = tf.paragraphs[0]
    p.level = level
    run = p.add_run()
    run.text = text
    _set_run_font(run, font_name)


def _fill_points(tf: Any, points: list[Any], font_name: str) -> None:
    """多段要点：clear + add_paragraph + level（首段复用已有段落）。"""
    tf.clear()
    for i, pt in enumerate(points):
        text = pt["text"] if isinstance(pt, dict) else pt
        level = pt.get("level", 0) if isinstance(pt, dict) else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        _set_run_font(run, font_name)


def _clear_slides(pres: Any) -> None:
    """清掉模板文件自带的示例页（只复用版式/主题；内置空模板为 no-op）。

    同时删除 presentation 到 slide 的 rel——否则孤立 slide part 仍会被
    写入输出包（zip 重复条目告警 + 文件膨胀）。
    """
    from pptx.oxml.ns import qn  # noqa: PLC0415
    xml_slides = pres.slides._sldIdLst
    for sldId in list(xml_slides):
        rId = sldId.get(qn("r:id"))
        if rId:
            pres.part.drop_rel(rId)
        xml_slides.remove(sldId)


def render_page(pres: Any, layout: Any, page: dict[str, Any], font_name: str,
                warnings: list[str]) -> None:
    """渲染单页：add_slide + 按角色填占位符；未匹配占位符不填（版式默认文本兜底）。"""
    slide = pres.slides.add_slide(layout)
    roles, _ = normalize._resolve_roles(list(slide.placeholders))
    title_ph = roles.get("title")
    points_ph = roles.get("points")
    caption_ph = roles.get("caption")
    image_ph = roles.get("image")

    content = page["content"]
    if page.get("title") and title_ph is not None:
        _fill_text_paragraph(title_ph.text_frame, page["title"], font_name)
    if content["points"] and points_ph is not None:
        _fill_points(points_ph.text_frame, content["points"], font_name)
    if content["image"]:
        if image_ph is not None:
            img_path = Path(content["image"])
            if not img_path.exists():
                raise FileNotFoundError(
                    f"pages[{page['index']}] 图片不存在: {img_path}"
                )
            image_ph.insert_picture(str(img_path))
        else:
            warnings.append(
                f"pages[{page['index']}] 声明了图片但版式无图片占位符——已跳过（版式默认内容兜底）"
            )
    if content["caption"]:
        target = caption_ph
        if target is None and not content["points"] and points_ph is not None:
            target = points_ph  # 图片页文本占位符即 caption（如 Picture with Caption 版式）
        if target is not None:
            _fill_text_paragraph(target.text_frame, content["caption"], font_name)
        else:
            warnings.append(
                f"pages[{page['index']}] 声明了 caption 但版式无文本占位符——已跳过"
            )


def render(envelope: dict[str, Any]) -> dict[str, Any]:
    """读信封渲染，返回结果 dict（供 CLI 打印 / 测试直调）。"""
    pages = envelope["pages"]
    out_path = Path(envelope["output"])
    font_name = envelope.get("font") or normalize.DEFAULT_FONT
    used_layouts = [p["layout"] for p in pages]

    # 渲染前 fail-fast：解析 + 硬合规校验（文件模板与内置默认模板同判据）
    deck_file, plan, issues = normalize.resolve_deck_layouts(used_layouts)
    if issues:
        raise ValueError(
            "模板硬合规未通过（渲染前 fail-fast）：\n"
            + "\n".join(f"  - {i}" for i in issues)
        )

    pptx = normalize._ensure_pptx()
    warnings: list[str] = []
    if deck_file is not None:
        # 模板形态诊断（页面层手工设计 / 未注册页面清单）——advisory
        warnings.extend(normalize.diagnose_template(normalize.REFERENCE_DIR / deck_file))
    pres = (
        pptx.Presentation()
        if deck_file is None
        else pptx.Presentation(str(normalize.REFERENCE_DIR / deck_file))
    )
    # 模板文件可能自带示例页——只复用版式/主题，输出不含模板自身 slide
    _clear_slides(pres)
    for page in pages:
        m_idx, idx, _kind = plan[page["layout"]]
        layout = pres.slide_masters[m_idx].slide_layouts[idx]
        render_page(pres, layout, page, font_name, warnings)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    pres.save(str(out_path))
    return {"pptx": str(out_path), "pages": len(pages), "warnings": warnings}


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="ppt_writer 渲染器（command 节点调用）")
    parser.add_argument("--envelope", required=True, help="翻译器写入的运行信封 JSON 路径")
    args = parser.parse_args(argv)
    try:
        normalize._ensure_pptx()
    except ImportError as e:
        print(f"错误: {e}", file=sys.stderr)
        return 2
    try:
        envelope = normalize.read_envelope_file(args.envelope)
        result = render(envelope)
    except (ValueError, FileNotFoundError, RuntimeError, KeyError) as e:
        print(f"渲染失败: {e}", file=sys.stderr)
        return 1
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
